"""
DJ Foundations Slide Rebuilder
Fixes broken Figma-to-PPTX export by repositioning elements based on reference images.

Reference images are 4000x2250px. Slide is 10.00" x 5.62".
Conversion: 400px = 1 inch (both directions)
"""

import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import copy

# Paths
BASE_DIR = Path(r"C:\Users\User\Documents\0000000\My Classes\DJ Foundations")
SOURCE_PPTX = BASE_DIR / "source material" / "DJ_Foundations_Styled (4).pptx"
OUTPUT_PPTX = BASE_DIR / "DJ_Foundations.pptx"
REFERENCE_DIR = BASE_DIR / "source material" / "images" / "slides"
BACKUP_DIR = BASE_DIR / "backups"

# Slide dimensions
SLIDE_WIDTH_INCHES = 10.0
SLIDE_HEIGHT_INCHES = 5.625  # 16:9 aspect ratio

# Reference image dimensions
REF_WIDTH_PX = 4000
REF_HEIGHT_PX = 2250

# Conversion factor: pixels to inches
PX_TO_INCH = SLIDE_WIDTH_INCHES / REF_WIDTH_PX  # 0.0025 or 400px = 1 inch

def px_to_inches(px):
    """Convert pixels to inches."""
    return px * PX_TO_INCH

def backup_current():
    """Create a backup of the current PPTX."""
    BACKUP_DIR.mkdir(exist_ok=True)
    import datetime
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = BACKUP_DIR / f"DJ_Foundations_backup_{ts}.pptx"
    if OUTPUT_PPTX.exists():
        shutil.copy(OUTPUT_PPTX, backup_path)
        print(f"Backed up to: {backup_path}")
    return backup_path

def reset_to_original():
    """Copy original export to working file."""
    if SOURCE_PPTX.exists():
        shutil.copy(SOURCE_PPTX, OUTPUT_PPTX)
        print(f"Reset to original: {SOURCE_PPTX.name}")
    else:
        print(f"ERROR: Source not found: {SOURCE_PPTX}")

def analyze_slide(prs, slide_num):
    """Analyze shapes on a slide and print their properties."""
    slide = prs.slides[slide_num - 1]
    print(f"\n{'='*60}")
    print(f"SLIDE {slide_num} ANALYSIS")
    print(f"{'='*60}")

    for i, shape in enumerate(slide.shapes):
        print(f"\nShape {i}: {shape.shape_type}")
        print(f"  Name: {shape.name}")
        print(f"  Position: left={shape.left.inches:.2f}\", top={shape.top.inches:.2f}\"")
        print(f"  Size: width={shape.width.inches:.2f}\", height={shape.height.inches:.2f}\"")

        if shape.has_text_frame:
            text = shape.text_frame.text[:100].replace('\n', ' ')
            print(f"  Text: {text}...")

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"  [PICTURE]")

def get_shape_by_text(slide, text_fragment):
    """Find a shape containing specific text."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            if text_fragment.lower() in shape.text_frame.text.lower():
                return shape
    return None

def get_shapes_by_type(slide, shape_type):
    """Get all shapes of a specific type."""
    return [s for s in slide.shapes if s.shape_type == shape_type]

def get_text_shapes(slide):
    """Get all shapes with text frames."""
    return [s for s in slide.shapes if s.has_text_frame]

def fix_text_box_height(shape, max_height_inches=4.0):
    """Fix text box height to be reasonable (not 4.95" default)."""
    if shape.height.inches > max_height_inches:
        # Estimate content height based on text
        text = shape.text_frame.text if shape.has_text_frame else ""
        lines = text.count('\n') + 1
        # Rough estimate: ~0.3" per line
        estimated_height = max(0.5, lines * 0.25)
        shape.height = Inches(min(estimated_height, max_height_inches))

def fix_slide_18(prs):
    """
    Slide 18: Two-Column Text Only (Take It Further)
    Reference measurements (from 4000x2250 image):
    - Title: centered top, ~60px top
    - Left column: ~80px left (0.2"), ~500px top (1.25"), ~1840px wide (4.6")
    - Right column: ~2080px left (5.2"), ~500px top (1.25"), ~1840px wide (4.6")
    """
    print("\nFIXING SLIDE 18: Two-column text layout")
    slide = prs.slides[17]  # 0-indexed

    text_shapes = get_text_shapes(slide)
    print(f"  Found {len(text_shapes)} text shapes")

    # Find the title and content shapes
    title_shape = None
    content_shapes = []

    for shape in text_shapes:
        text = shape.text_frame.text.lower()
        if "take it further" in text or shape.top.inches < 0.8:
            title_shape = shape
        else:
            content_shapes.append(shape)

    # Identify columns by content keywords
    left_content = None
    right_content = None

    for shape in content_shapes:
        text = shape.text_frame.text
        # Left column has "Slam Academy", "History of DJing", etc.
        if "Slam Academy" in text or "History of DJing" in text:
            left_content = shape
        # Right column has "Manual beatmatching", "Harmonic mixing", etc.
        elif "Manual beatmatching" in text or "Harmonic mixing" in text:
            right_content = shape

    # If we couldn't identify by content, use current position
    if not left_content and not right_content and len(content_shapes) >= 2:
        sorted_shapes = sorted(content_shapes, key=lambda s: s.left.inches)
        left_content = sorted_shapes[0]
        right_content = sorted_shapes[1] if len(sorted_shapes) > 1 else None

    # Apply fixed positions based on reference image measurements
    if title_shape:
        title_shape.left = Inches(0.2)
        title_shape.top = Inches(0.15)
        title_shape.width = Inches(9.6)
        title_shape.height = Inches(0.9)
        print(f"  Title positioned at (0.2\", 0.15\")")

    if left_content:
        left_content.left = Inches(0.2)
        left_content.top = Inches(1.25)
        left_content.width = Inches(4.6)
        left_content.height = Inches(4.0)
        print(f"  Left column at (0.2\", 1.25\"), width=4.6\"")

    if right_content:
        right_content.left = Inches(5.2)
        right_content.top = Inches(1.25)
        right_content.width = Inches(4.6)
        right_content.height = Inches(4.0)
        print(f"  Right column at (5.2\", 1.25\"), width=4.6\"")

def fix_slide_9(prs):
    """
    Slide 9: Beats/Bars/Phrases - Special Layout
    Reference measurements (from 4000x2250 image):
    - Title: ~60px top, spans left portion
    - Text content: ~80px left (0.2"), ~400px top (1.0"), ~1600px wide (4.0")
    - Song-body diagram (waveform+beat/bar/phrase labels): right side
      ~1680px left (4.2"), ~380px top (0.95"), ~2000px wide (5.0"), ~880px tall (2.2")
    - 16-beat diagram: bottom center
      ~320px left (0.8"), ~1720px top (4.3%), ~3360px wide (8.4"), ~400px tall (1.0")
    """
    print("\nFIXING SLIDE 9: Beats/Bars/Phrases diagram layout")
    slide = prs.slides[8]  # 0-indexed

    # Get all shapes
    pictures = get_shapes_by_type(slide, MSO_SHAPE_TYPE.PICTURE)
    text_shapes = get_text_shapes(slide)

    print(f"  Found {len(pictures)} pictures, {len(text_shapes)} text shapes")

    # Find title shape and content shapes
    title_shape = None
    content_shapes = []

    for shape in text_shapes:
        if shape.top.inches < 0.8:
            title_shape = shape
        else:
            content_shapes.append(shape)

    # Position title - spans left side
    if title_shape:
        title_shape.left = Inches(0.2)
        title_shape.top = Inches(0.15)
        title_shape.width = Inches(9.6)
        title_shape.height = Inches(0.8)
        print(f"  Title at (0.2\", 0.15\")")

    # Position content text on left - CONSTRAINED WIDTH
    for shape in content_shapes:
        shape.left = Inches(0.2)
        shape.top = Inches(1.0)
        shape.width = Inches(4.0)  # Key fix: constrain to left column only
        shape.height = Inches(2.8)  # Enough for bullet points
        print(f"  Content text at (0.2\", 1.0\"), width=4.0\"")

    # Position diagrams based on reference
    if len(pictures) >= 2:
        # Identify diagrams by aspect ratio
        # Song-body diagram is more square (~2.3 aspect), 16-beat is very wide (~5.8 aspect)
        pic1, pic2 = pictures[0], pictures[1]

        ar1 = pic1.width.inches / pic1.height.inches if pic1.height.inches > 0 else 1
        ar2 = pic2.width.inches / pic2.height.inches if pic2.height.inches > 0 else 1

        # The 16-beat diagram is much wider (higher aspect ratio)
        if ar1 > ar2:
            beat_diagram = pic1
            song_body = pic2
        else:
            beat_diagram = pic2
            song_body = pic1

        # Position song-body diagram (waveform with beat/bar/phrase labels): top-right
        song_body.left = Inches(4.2)
        song_body.top = Inches(0.95)
        song_body.width = Inches(5.6)
        song_body.height = Inches(2.2)
        print(f"  Song-body diagram at (4.2\", 0.95\"), size=5.6\"x2.2\"")

        # Position 16-beat diagram: bottom center
        beat_diagram.left = Inches(0.8)
        beat_diagram.top = Inches(4.0)
        beat_diagram.width = Inches(8.4)
        beat_diagram.height = Inches(1.2)
        print(f"  16-beat diagram at (0.8\", 4.0\"), size=8.4\"x1.2\"")

    elif len(pictures) == 1:
        pic = pictures[0]
        pic.left = Inches(4.2)
        pic.top = Inches(0.95)
        pic.width = Inches(5.6)
        pic.height = Inches(3.5)
        print(f"  Single diagram at (4.2\", 0.95\")")

def fix_slide_2(prs):
    """
    Slide 2: Who Am I? - text + meme character + DJ photo
    Reference measurements (from 4000x2250 image):
    - Title: centered top, ~60px top (0.15")
    - Text content: ~80px left (0.2"), ~440px top (1.1"), ~1760px wide (4.4")
    - Meme character: bottom left, ~120px left (0.3"), ~1200px top (3.0"), ~800x880px (2.0"x2.2")
    - DJ photo: right side, ~2000px left (5.0"), ~320px top (0.8"), fills to edge

    IMPORTANT: Preserve aspect ratios to avoid squishing!
    """
    print("\nFIXING SLIDE 2: Who Am I? layout")
    slide = prs.slides[1]  # 0-indexed

    pictures = get_shapes_by_type(slide, MSO_SHAPE_TYPE.PICTURE)
    text_shapes = get_text_shapes(slide)

    print(f"  Found {len(pictures)} pictures, {len(text_shapes)} text shapes")

    # Position title - don't let images overlap it
    title_shape = None
    content_shape = None

    for shape in text_shapes:
        if shape.top.inches < 0.8:
            title_shape = shape
        else:
            content_shape = shape

    if title_shape:
        title_shape.left = Inches(0.2)
        title_shape.top = Inches(0.15)
        title_shape.width = Inches(9.6)
        title_shape.height = Inches(0.8)
        print(f"  Title at (0.2\", 0.15\")")

    if content_shape:
        content_shape.left = Inches(0.2)
        content_shape.top = Inches(1.1)
        content_shape.width = Inches(4.4)
        content_shape.height = Inches(1.8)
        print(f"  Content text at (0.2\", 1.1\"), width=4.4\"")

    # Position images - PRESERVE ASPECT RATIOS
    if len(pictures) >= 2:
        pic1, pic2 = pictures[0], pictures[1]

        # Identify by original aspect ratio
        # Meme character is roughly square, DJ photo is also roughly square but larger
        area1 = pic1.width.inches * pic1.height.inches
        area2 = pic2.width.inches * pic2.height.inches

        if area1 > area2:
            dj_photo = pic1
            meme_char = pic2
        else:
            dj_photo = pic2
            meme_char = pic1

        # Get original aspect ratios
        dj_aspect = dj_photo.width.inches / dj_photo.height.inches if dj_photo.height.inches > 0 else 1
        meme_aspect = meme_char.width.inches / meme_char.height.inches if meme_char.height.inches > 0 else 1

        # Position DJ photo: right side, below title, preserve aspect ratio
        # Target: 5.0" wide, calculate height from aspect ratio
        dj_target_width = 5.0
        dj_target_height = dj_target_width / dj_aspect
        # Cap height to not exceed slide
        dj_target_height = min(dj_target_height, 4.8)
        # Recalculate width if height was capped
        if dj_target_height == 4.8:
            dj_target_width = dj_target_height * dj_aspect

        dj_photo.left = Inches(5.0)
        dj_photo.top = Inches(0.8)  # Below title
        dj_photo.width = Inches(dj_target_width)
        dj_photo.height = Inches(dj_target_height)
        print(f"  DJ photo at (5.0\", 0.8\"), size={dj_target_width:.1f}\"x{dj_target_height:.1f}\" (aspect preserved)")

        # Position meme character: bottom-left, preserve aspect ratio
        # Target: 2.0" wide, calculate height from aspect ratio
        meme_target_width = 2.0
        meme_target_height = meme_target_width / meme_aspect
        # Cap to reasonable size
        meme_target_height = min(meme_target_height, 2.5)
        if meme_target_height == 2.5:
            meme_target_width = meme_target_height * meme_aspect

        meme_char.left = Inches(0.3)
        meme_char.top = Inches(3.0)
        meme_char.width = Inches(meme_target_width)
        meme_char.height = Inches(meme_target_height)
        print(f"  Meme character at (0.3\", 3.0\"), size={meme_target_width:.1f}\"x{meme_target_height:.1f}\" (aspect preserved)")

    elif len(pictures) == 1:
        pic = pictures[0]
        aspect = pic.width.inches / pic.height.inches if pic.height.inches > 0 else 1
        target_width = 5.0
        target_height = target_width / aspect
        pic.left = Inches(5.0)
        pic.top = Inches(0.8)
        pic.width = Inches(target_width)
        pic.height = Inches(min(target_height, 4.8))
        print(f"  Single image at (5.0\", 0.8\")")

def fix_two_column_slide(prs, slide_num, description=""):
    """
    Generic fix for two-column text+image slides.
    Reference measurements (from 4000x2250px images, 400px=1"):
    - Title: full width centered, ~60px top (0.15")
    - Text content: ~80px left (0.2"), ~440px top (1.1"), ~1760px wide (4.4")
    - Image: starts ~1800px from left (4.5"), extends to right edge
    """
    print(f"\nFIXING SLIDE {slide_num}: {description}")
    slide = prs.slides[slide_num - 1]

    pictures = get_shapes_by_type(slide, MSO_SHAPE_TYPE.PICTURE)
    text_shapes = get_text_shapes(slide)

    # Find title vs content
    title_shape = None
    content_shapes = []

    for shape in text_shapes:
        if shape.top.inches < 0.8:
            title_shape = shape
        else:
            content_shapes.append(shape)

    # Position title - match reference position exactly
    if title_shape:
        title_shape.left = Inches(0.0)
        title_shape.top = Inches(0.15)
        title_shape.width = Inches(10.0)
        title_shape.height = Inches(0.85)

    # Position content text on left
    for shape in content_shapes:
        shape.left = Inches(0.2)
        shape.top = Inches(1.1)
        shape.width = Inches(4.4)
        shape.height = Inches(2.2)

    # Position image on right - extend to edge
    if pictures:
        pic = pictures[0]
        pic.left = Inches(4.5)
        pic.top = Inches(0.9)
        pic.width = Inches(5.5)
        pic.height = Inches(4.7)

def fix_slide_12(prs):
    """
    Slide 12: Exporting to USB - special layout
    - Title at top
    - Text on left (bullet points)
    - USB drive image on right
    - Screenshot at bottom left
    """
    print("\nFIXING SLIDE 12: Exporting to USB layout")
    slide = prs.slides[11]

    pictures = get_shapes_by_type(slide, MSO_SHAPE_TYPE.PICTURE)
    text_shapes = get_text_shapes(slide)

    # Find title and content
    title_shape = None
    content_shape = None

    for shape in text_shapes:
        if shape.top.inches < 0.8:
            title_shape = shape
        elif shape.width.inches > 2:  # Main content, not a label
            content_shape = shape

    if title_shape:
        title_shape.left = Inches(0.0)
        title_shape.top = Inches(0.15)
        title_shape.width = Inches(10.0)
        title_shape.height = Inches(0.85)

    if content_shape:
        content_shape.left = Inches(0.2)
        content_shape.top = Inches(1.2)
        content_shape.width = Inches(5.5)
        content_shape.height = Inches(2.5)

    # Position images - USB drive on right, screenshot at bottom
    if len(pictures) >= 2:
        # Sort by size to identify USB (larger) vs screenshot (smaller/wider)
        sorted_pics = sorted(pictures, key=lambda p: p.width.inches * p.height.inches, reverse=True)
        usb_pic = sorted_pics[0]
        screenshot = sorted_pics[1] if len(sorted_pics) > 1 else None

        # USB drive - right side
        usb_pic.left = Inches(5.5)
        usb_pic.top = Inches(1.0)
        usb_pic.width = Inches(4.4)
        usb_pic.height = Inches(4.4)

        # Screenshot - bottom left
        if screenshot:
            screenshot.left = Inches(0.9)
            screenshot.top = Inches(3.8)
            screenshot.width = Inches(4.6)
            screenshot.height = Inches(0.75)

def fix_slide_17(prs):
    """
    Slide 17: Practice & Next Steps - two columns of text above large image
    - Title at top
    - Two text columns side by side
    - Large DJ equipment image spanning bottom
    """
    print("\nFIXING SLIDE 17: Practice & Next Steps layout")
    slide = prs.slides[16]

    pictures = get_shapes_by_type(slide, MSO_SHAPE_TYPE.PICTURE)
    text_shapes = get_text_shapes(slide)

    # Find title and content shapes
    title_shape = None
    content_shapes = []

    for shape in text_shapes:
        if shape.top.inches < 0.6:
            title_shape = shape
        else:
            content_shapes.append(shape)

    if title_shape:
        title_shape.left = Inches(0.0)
        title_shape.top = Inches(0.15)
        title_shape.width = Inches(10.0)
        title_shape.height = Inches(0.75)

    # Position the two content columns
    if len(content_shapes) >= 2:
        # Sort by left position
        sorted_content = sorted(content_shapes, key=lambda s: s.left.inches)
        left_col = sorted_content[0]
        right_col = sorted_content[1]

        left_col.left = Inches(0.5)
        left_col.top = Inches(1.0)
        left_col.width = Inches(4.0)
        left_col.height = Inches(1.0)

        right_col.left = Inches(5.0)
        right_col.top = Inches(1.0)
        right_col.width = Inches(4.5)
        right_col.height = Inches(1.0)
    elif len(content_shapes) == 1:
        # Single content shape
        content_shapes[0].left = Inches(0.5)
        content_shapes[0].top = Inches(1.0)
        content_shapes[0].width = Inches(9.0)
        content_shapes[0].height = Inches(1.2)

    # Position the DJ equipment image - full width at bottom
    if pictures:
        pic = pictures[0]
        pic.left = Inches(0.5)
        pic.top = Inches(2.0)
        pic.width = Inches(9.0)
        pic.height = Inches(3.5)

def fix_all_text_heights(prs):
    """Fix all text box heights across all slides."""
    print("\nFIXING ALL TEXT BOX HEIGHTS")
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.height.inches > 4.5:
                old_height = shape.height.inches
                fix_text_box_height(shape, 4.0)
                if old_height != shape.height.inches:
                    print(f"  Slide {i+1}: Fixed height {old_height:.2f}\" -> {shape.height.inches:.2f}\"")

def main():
    """Main rebuild process."""
    print("=" * 60)
    print("DJ FOUNDATIONS SLIDE REBUILDER")
    print("=" * 60)

    # Backup current
    backup_current()

    # Reset to original export
    reset_to_original()

    # Load presentation
    print(f"\nLoading: {OUTPUT_PPTX}")
    prs = Presentation(str(OUTPUT_PPTX))

    print(f"Slides: {len(prs.slides)}")
    print(f"Slide size: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")

    # Fix only the truly broken slides
    fix_slide_18(prs)  # 76.3% - overlapping columns (CRITICAL)
    fix_slide_9(prs)   # 87.8% - diagram positioning (CRITICAL)
    fix_slide_2(prs)   # 81.4% - needs meme + DJ photo positioning

    # Fix text heights only (don't move shapes on slides that are already OK)
    fix_all_text_heights(prs)

    # Save
    print(f"\nSaving: {OUTPUT_PPTX}")
    prs.save(str(OUTPUT_PPTX))
    print("Done!")

    print("\n" + "=" * 60)
    print("NEXT STEPS:")
    print("1. Run verify_slides.py to check results")
    print("2. Review slide_diffs/ for visual comparison")
    print("3. Iterate on problem slides")
    print("=" * 60)

if __name__ == "__main__":
    import sys

    if len(sys.argv) > 1 and sys.argv[1] == "analyze":
        # Analysis mode - show slide structure
        prs = Presentation(str(OUTPUT_PPTX))
        for slide_num in [9, 18, 2]:  # Analyze problem slides
            analyze_slide(prs, slide_num)
    else:
        main()
